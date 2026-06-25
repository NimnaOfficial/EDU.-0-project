import io
import logging
import asyncio
import copy
from pypdf import PdfWriter, PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

class DocumentMerger:
    """
    Advanced Infrastructure logic for stitching multiple documents.
    Operates entirely in RAM using asynchronous threading for zero lag.
    """

    # ==========================================
    # PDF ENGINE
    # ==========================================
    @staticmethod
    async def merge_pdfs(pdf_buffers: list[io.BytesIO]) -> io.BytesIO:
        try:
            return await asyncio.to_thread(DocumentMerger._merge_pdf_sync, pdf_buffers)
        except Exception as e:
            logger.error(f"Failed to execute PDF merge: {str(e)}")
            raise

    @staticmethod
    def _merge_pdf_sync(pdf_buffers: list[io.BytesIO]) -> io.BytesIO:
        logger.info(f"Initiating PDF Merge sequence for {len(pdf_buffers)} documents.")
        merger = PdfWriter()
        for pdf_buffer in pdf_buffers:
            pdf_buffer.seek(0)
            merger.append(PdfReader(pdf_buffer))
        output_buffer = io.BytesIO()
        merger.write(output_buffer)
        output_buffer.seek(0)
        return output_buffer

    # ==========================================
    # PPTX ENGINE (Patched for Image Rendering)
    # ==========================================
    @staticmethod
    async def merge_pptxs(pptx_buffers: list[io.BytesIO]) -> io.BytesIO:
        try:
            return await asyncio.to_thread(DocumentMerger._merge_pptx_sync, pptx_buffers)
        except Exception as e:
            logger.error(f"Failed to execute PPTX merge: {str(e)}")
            raise

    @staticmethod
    def _merge_pptx_sync(buffers: list[io.BytesIO]) -> io.BytesIO:
        logger.info(f"Initiating PPTX Merge sequence for {len(buffers)} documents.")
        
        buffers[0].seek(0)
        master_prs = Presentation(buffers[0])
        
        try:
            blank_layout = master_prs.slide_layouts[6]
        except IndexError:
            blank_layout = master_prs.slide_layouts[0] 
        
        for buffer in buffers[1:]:
            buffer.seek(0)
            sub_prs = Presentation(buffer)
            
            for slide in sub_prs.slides:
                new_slide = master_prs.slides.add_slide(blank_layout)
                
                for shape in slide.shapes:
                    try:
                        # SMART PARSING: Handle Pictures Properly
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            # 1. Extract the raw image binary from the old slide
                            image_bytes = shape.image.blob
                            image_stream = io.BytesIO(image_bytes)
                            
                            # 2. Re-embed it physically into the new master presentation
                            new_slide.shapes.add_picture(
                                image_stream, 
                                shape.left, shape.top, shape.width, shape.height
                            )
                        else:
                            # SAFE COPY: For text boxes and basic shapes without media links
                            new_el = copy.deepcopy(shape.element)
                            new_slide.shapes._spTree.append(new_el)
                            
                    except Exception as e:
                        logger.warning(f"Skipped complex shape during merge: {e}")
                        
        output_buffer = io.BytesIO()
        master_prs.save(output_buffer)
        output_buffer.seek(0)
        
        return output_buffer