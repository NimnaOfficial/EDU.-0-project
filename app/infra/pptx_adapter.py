import io
import logging
import asyncio
import html
from html.parser import HTMLParser
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from typing import Any

from domain.parser import ParsedH5P

logger = logging.getLogger(__name__)

# ==========================================
# ADVANCED RICH TEXT ENGINE
# ==========================================
class H5PRichTextParser(HTMLParser):
    """
    A lightweight, high-performance HTML parser that translates web 
    typography directly into Microsoft PowerPoint XML runs.
    """
    def __init__(self, paragraph):
        super().__init__()
        self.paragraph = paragraph
        self.bold = False
        self.italic = False
        self.underline = False

    def handle_starttag(self, tag, attrs):
        if tag in ['b', 'strong']: self.bold = True
        elif tag in ['i', 'em']: self.italic = True
        elif tag in ['u']: self.underline = True
        elif tag in ['br', 'p']:
            run = self.paragraph.add_run()
            run.text = "\n"

    def handle_endtag(self, tag):
        if tag in ['b', 'strong']: self.bold = False
        elif tag in ['i', 'em']: self.italic = False
        elif tag in ['u']: self.underline = False

    def handle_data(self, data):
        text = html.unescape(data).strip('\n')
        if not text:
            return
            
        run = self.paragraph.add_run()
        run.text = text
        run.font.bold = self.bold
        run.font.italic = self.italic
        run.font.underline = self.underline
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(44, 62, 80) # Modern UI Dark Grey


# ==========================================
# PPTX BUILDER INFRASTRUCTURE
# ==========================================
class PPTXBuilder:
    """
    Advanced Infrastructure logic for generating PowerPoint files.
    Constructs 16:9 widescreen slides, parses rich text, and handles Z-Index layering.
    """

    # 16:9 Widescreen Configuration
    SLIDE_WIDTH_INCHES = 13.333
    SLIDE_HEIGHT_INCHES = 7.5

    @staticmethod
    async def build_presentation(h5p_data: ParsedH5P) -> io.BytesIO:
        """Runs the compilation asynchronously to prevent blocking the main event loop."""
        try:
            return await asyncio.to_thread(PPTXBuilder._compile_sync, h5p_data)
        except Exception as e:
            logger.error(f"Failed to build PPTX infrastructure: {str(e)}")
            raise

    @staticmethod
    def _compile_sync(data: ParsedH5P) -> io.BytesIO:
        logger.info(f"Initiating V2.0 PPTX build sequence for: {data.metadata.title}")
        
        prs = Presentation()
        
        # 1. Force Modern 16:9 Widescreen Aspect Ratio
        prs.slide_width = Inches(PPTXBuilder.SLIDE_WIDTH_INCHES)
        prs.slide_height = Inches(PPTXBuilder.SLIDE_HEIGHT_INCHES)
        
        title_slide_layout = prs.slide_layouts[0] 
        blank_slide_layout = prs.slide_layouts[6] 
        
        # 2. Build the Master Title Slide
        title_slide = prs.slides.add_slide(title_slide_layout)
        title: Any = title_slide.shapes.title
        subtitle: Any = title_slide.placeholders[1]

        if title and title.has_text_frame:
            title.text_frame.text = data.metadata.title
            for paragraph in title.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(48)
                paragraph.font.color.rgb = RGBColor(41, 128, 185) 

        if subtitle and subtitle.has_text_frame:
            subtitle.text_frame.text = f"Extracted from: {data.metadata.mainLibrary}\nCompiled by EDU. 0 Engine V2.0"

        # 3. Process Content Slides
        slide_array = data.content.presentation.get('slides', [])
        
        if not slide_array:
             slide_array = [{"elements": [{"action": {"params": {"text": "Notice: Layout requires advanced parsing."}}}]}]

        for idx, slide_data in enumerate(slide_array):
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Extract elements and sort them by Z-Index (Painter's Algorithm)
            elements = slide_data.get('elements', [])
            
            # Pylance strict type-checking fix: Guarantee an integer return
            def _z_index_key(x: dict[str, Any]) -> int:
                z_val = x.get('z')
                if isinstance(z_val, (int, float)):
                    return int(z_val)
                if isinstance(z_val, str):
                    try:
                        return int(float(z_val))
                    except ValueError:
                        return 0
                return 0

            elements.sort(key=_z_index_key)
            
            for element in elements:
                PPTXBuilder._map_element_to_slide(slide, element, data.raw_assets)

        # 4. Save to RAM safely
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        logger.info("PPTX build sequence complete. Returning binary stream.")
        return pptx_buffer

    @staticmethod
    def _map_element_to_slide(slide, element: dict, raw_assets: dict):
        """Maps H5P JSON to perfectly aligned PPTX shapes."""
        try:
            # Safely grab percentage coordinates
            x_pct = element.get('x', 10) / 100
            y_pct = element.get('y', 20) / 100
            w_pct = element.get('width', 80) / 100
            h_pct = element.get('height', 10) / 100

            # Advanced UI Math: Translate percentages to 16:9 widescreen physical inches
            left = Inches(PPTXBuilder.SLIDE_WIDTH_INCHES * x_pct)
            top = Inches(PPTXBuilder.SLIDE_HEIGHT_INCHES * y_pct)
            width = Inches(PPTXBuilder.SLIDE_WIDTH_INCHES * w_pct)
            height = Inches(PPTXBuilder.SLIDE_HEIGHT_INCHES * h_pct)

            action = element.get('action', {})
            library = action.get('library', '').split(' ')[0] 
            params = action.get('params', {})

            # --- ROUTE A: RICH TEXT ELEMENTS ---
            if 'Text' in library:
                raw_html = params.get('text', '')
                
                # Create text box
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                
                # Execute custom HTML-to-XML Parsing
                p = tf.add_paragraph()
                parser = H5PRichTextParser(p)
                parser.feed(raw_html)

            # --- ROUTE B: IMAGE ELEMENTS ---
            elif 'Image' in library:
                image_path = params.get('file', {}).get('path', '')
                
                if image_path in raw_assets:
                    image_stream = io.BytesIO(raw_assets[image_path])
                    slide.shapes.add_picture(image_stream, left, top, width, height)

        except Exception as e:
            logger.warning(f"Engine skipped an overlapping or corrupted shape element: {e}")